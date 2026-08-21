import yaml
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData

class PptxRenderer:
    def __init__(self, template_path: str, yaml_path: str, output_path: str):
        self.template_path = template_path
        self.yaml_path = yaml_path
        self.output_path = output_path

    def render(self):
        prs = Presentation(self.template_path)
        with open(self.yaml_path, 'r', encoding='utf-8') as f:
            yaml_data = yaml.safe_load(f)
                    
        output_slides = yaml_data.get('output_slides', {})

        for idx in range(len(output_slides)):
            print(f"-> Updating Slide {idx}...")
            self._update_slide(prs.slides[idx], output_slides[idx])
            
        prs.save(self.output_path)
        print(f"\n[Success] Saved at: {self.output_path}")

    def _update_slide(self, slide, slide_data):
        elements = slide_data.get('elements', [])
        
        text_shapes = [s for s in slide.shapes if s.has_text_frame]
        table_shapes = [s for s in slide.shapes if s.has_table]
        chart_shapes = [s for s in slide.shapes if s.has_chart]
        
        yaml_texts = [e for e in elements if e.get('role') in ['slide-title', 'body-text', 'text', 'caption']]
        yaml_tables = [e for e in elements if e.get('role') == 'table']
        yaml_charts = [e for e in elements if e.get('role') in ['chart-bar', 'chart-line']]
        
        for i, text_element in enumerate(yaml_texts):
            if i < len(text_shapes):
                new_text = text_element.get('text', "")
                if new_text:
                    self._replace_text_preserve_format(text_shapes[i], new_text)
                    
        for i, table_element in enumerate(yaml_tables):
            if i < len(table_shapes):
                # self._update_table(table_shapes[i].table, table_element.get('data', []))
                self._update_table(slide, table_shapes[i], table_element.get('data', []))
                
        for i, chart_element in enumerate(yaml_charts):
            if i < len(chart_shapes):
                self._update_chart(chart_shapes[i].chart, chart_element.get('data', []))

    def _replace_text_preserve_format(self, shape, new_text):
        text_frame = shape.text_frame
        lines = str(new_text).split('\n')
        
        for i, line in enumerate(lines):
            if i < len(text_frame.paragraphs):
                p = text_frame.paragraphs[i]
                if p.runs:
                    p.runs[0].text = line
                    for r in p.runs[1:]:
                        r.text = ""
                else:
                    p.text = line
            else:
                p = text_frame.add_paragraph()
                p.text = line
                
        if len(text_frame.paragraphs) > len(lines):
            for i in range(len(text_frame.paragraphs) - 1, len(lines) - 1, -1):
                p = text_frame.paragraphs[i]
                p_element = p._element
                p_element.getparent().remove(p_element)

    def _update_table(self, slide, table_shape, data_list):            
        df = pd.DataFrame(data_list)
        columns = df.columns.tolist()
        
        target_rows = len(df) + 1
        target_cols = len(columns)
        
        current_rows = len(table_shape.table.rows)
        current_cols = len(table_shape.table.columns)
        
        if target_rows == current_rows and target_cols == current_cols:
            table = table_shape.table
            for row_idx in range(target_rows):
                for col_idx in range(target_cols):
                    cell = table.cell(row_idx, col_idx)
                    if row_idx == 0:
                        cell.text = str(columns[col_idx])
                    else:
                        cell.text = str(df.iloc[row_idx - 1, col_idx])
        else:
            left = table_shape.left
            top = table_shape.top
            width = table_shape.width
            height = table_shape.height
            
            sp = table_shape._element
            sp.getparent().remove(sp)
            
            new_table_shape = slide.shapes.add_table(target_rows, target_cols, left, top, width, height)
            table = new_table_shape.table
            
            for row_idx in range(target_rows):
                for col_idx in range(target_cols):
                    cell = table.cell(row_idx, col_idx)
                    if row_idx == 0:
                        cell.text = str(columns[col_idx])
                    else:
                        cell.text = str(df.iloc[row_idx - 1, col_idx])

    def _update_chart(self, chart, data_list):
        if not data_list: 
            return
            
        chart_data = CategoryChartData()
        df = pd.DataFrame(data_list)
        
        categories = df.iloc[:, 0].astype(str).tolist()
        chart_data.categories = categories
        
        for col in df.columns[1:]:
            series_name = str(col)
            series_values = pd.to_numeric(df[col], errors='coerce').fillna(0).tolist()
            chart_data.add_series(series_name, series_values)
            
        chart.replace_data(chart_data)


if __name__ == "__main__":
    TEMPLATE_PPTX = "slides/slide5.pptx" 
    YAML_GENERATED = "slides/slide5_generated_doc.yaml"
    OUTPUT_PPTX = "slides/slide5_fix.pptx"
    
    renderer = PptxRenderer(
        template_path=TEMPLATE_PPTX,
        yaml_path=YAML_GENERATED,
        output_path=OUTPUT_PPTX
    )
    renderer.render()