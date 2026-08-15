import ast
import re
import pandas as pd
import yaml
import os
import shutil
import io
import base64
from pathlib import Path
from config import config
from langchain_core.prompts import ChatPromptTemplate
from langchain_openai import ChatOpenAI
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE
from math import sqrt, hypot
from typing import Dict, List, Any, Tuple
from pptxtopdf import convert
from pdf2image import convert_from_path

from file_utils import load_prompt_from_file
from pptx_analyser import _call_vision_model_v2

def emu_to_cm(emu: float) -> float:
    return round(emu / 360000.0, 2)

def get_shape_layout(shape) -> Dict[str, float]:
    return {
        "x": emu_to_cm(shape.left),
        "y": emu_to_cm(shape.top),
        "width": emu_to_cm(shape.width),
        "height": emu_to_cm(shape.height),
    }


def get_shape_center(shape) -> Tuple[float, float]:
    return shape.left + shape.width / 2, shape.top + shape.height / 2


def table_shape_to_df(shape) -> "pd.DataFrame | None":
    if not hasattr(shape, "table"):
        return None
    table = shape.table
    rows = table.rows
    cols = table.columns

    data = []
    for r in range(len(rows)):
        row_vals = []
        for c in range(len(cols)):
            cell = table.cell(r, c)
            txt = cell.text_frame.text if cell.text_frame else ""
            row_vals.append(txt.strip())
        data.append(row_vals)

    if data:
        header = data[0]
        body = data[1:] if len(data) > 1 else []
        if len(set(h or f"col_{i}" for i, h in enumerate(header))) == len(header) and any(h.strip() for h in header):
            df = pd.DataFrame(body, columns=[h if h else f"col_{i}" for i, h in enumerate(header)])
        else:
            df = pd.DataFrame(data)
    else:
        df = pd.DataFrame()
    return df


def chart_shape_to_df(shape) -> "pd.DataFrame | None":
    if not hasattr(shape, "chart"):
        return None
    chart = shape.chart

    categories = []
    try:
        if chart.plots and chart.plots[0].categories is not None:
            for cat in chart.plots[0].categories:
                if hasattr(cat, "label"):
                    categories.append(str(cat.label))
                else:
                    categories.append(str(cat))
        else:
            categories = None
    except Exception:
        categories = None

    series_data = {}
    max_len = 0
    for s in chart.series:
        name = s.name if s.name is not None else f"series_{len(series_data)}"
        values = []
        for v in (s.values or []):
            if hasattr(v, "value"):
                values.append(v.value)
            else:
                try:
                    values.append(v)
                except Exception:
                    values.append(v)
        series_data[str(name)] = values
        max_len = max(max_len, len(values))

    if categories is None:
        categories = [f"cat_{i + 1}" for i in range(max_len)]
    else:
        if len(categories) < max_len:
            categories = categories + [f"cat_{i + 1}" for i in range(len(categories), max_len)]
        elif len(categories) > max_len:
            categories = categories[:max_len]

    df = pd.DataFrame({"category": categories})
    for series_name, vals in series_data.items():
        if len(vals) < max_len:
            vals = vals + [None] * (max_len - len(vals))
        elif len(vals) > max_len:
            vals = vals[:max_len]
        df[series_name] = vals
    return df


def _convert_ppt_to_image(ppt_path: str, slide_number: int, output_folder_path: str):
    file_name = os.path.splitext(os.path.basename(ppt_path))[0]
    file_path = Path(ppt_path).parent
    temp_pdf_output_path = os.path.join(file_path, f"{file_name}_temp_pdf")
    temp_image_output_path = os.path.join(file_path, f"{file_name}_temp_images")
    os.makedirs(temp_pdf_output_path, exist_ok=True)
    os.makedirs(temp_image_output_path, exist_ok=True)

    try:
        convert(input_path=ppt_path, output_folder_path=temp_pdf_output_path)
        pdf_file_path = os.path.join(temp_pdf_output_path, f"{file_name}.pdf")
        if not os.path.exists(pdf_file_path):
            raise FileNotFoundError(f"PDF conversion failed: {pdf_file_path}")
        images = convert_from_path(
            pdf_file_path,
            # poppler_path=r"C:\poppler-25.07.0\Library\bin",
            poppler_path=r"C:\Users\hp450\anaconda3\envs\slide\Library\bin",
            fmt="jpeg",
            dpi=200
        )
        if not (0 <= slide_number - 1 < len(images)):
            return None, None, None, None
        target_slide_pil_image = images[slide_number - 1]
        width_px, height_px = target_slide_pil_image.size
        return target_slide_pil_image, width_px, height_px, temp_image_output_path
    except Exception as e:
        print(f"\n[CRITICAL ERROR] _convert_ppt_to_image: {e}\n")
        return None, None, None, None
    finally:
        if os.path.exists(temp_pdf_output_path):
            shutil.rmtree(temp_pdf_output_path)

class PptxParser:
    def __init__(self, pptx_path: Path):
        if not pptx_path.exists():
            raise FileNotFoundError(f"PPT file not found: {pptx_path}")
        self.presentation = Presentation(pptx_path)
        self.file_name = pptx_path.stem
        self.file_path = str(pptx_path)
        self.slide_count = len(self.presentation.slides)

        self.model = ChatOpenAI(
            base_url=config.BASE_URL,
            api_key=config.API_KEY,
            temperature=0,
            model=config.MODEL_NAME
        )
        self.table_information_extraction_prompt_template = self._create_table_information_extraction_prompt_template()

    def _create_table_information_extraction_prompt_template(self) -> ChatPromptTemplate:
        system_prompt = load_prompt_from_file("function_logic_extraction_prompt.txt")
        return ChatPromptTemplate.from_messages([
            ("system", system_prompt),
            ("human", """
                     table_caption:{table_caption}
                     table_data:{table_data}
                     """)
        ])

    def _get_shape_type(self, shape) -> str:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "table"

        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            if hasattr(shape, "chart"):
                chart_type = shape.chart.chart_type
                if chart_type in (
                        XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.COLUMN_STACKED, XL_CHART_TYPE.COLUMN_STACKED_100,
                        XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.BAR_STACKED, XL_CHART_TYPE.BAR_STACKED_100
                ):
                    return "chart-bar"
                elif chart_type in (
                        XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS, XL_CHART_TYPE.LINE_STACKED,
                        XL_CHART_TYPE.LINE_STACKED_100, XL_CHART_TYPE.LINE_MARKERS_STACKED,
                        XL_CHART_TYPE.LINE_MARKERS_STACKED_100
                ):
                    return "chart-line"
                else:
                    return "chart-other"
            return "chart"

        if shape.has_text_frame and shape.text.strip():
            return "text"

        return "other"

    def _extract_pptx_elements1(self, slide):
        content_elements = []
        for idx, shape in enumerate(slide.shapes):

            shape_type = self._get_shape_type(shape)
            if shape_type == "table":
                df = table_shape_to_df(shape)
                try:
                    content_elements.append(df.to_dict(orient='records'))
                except Exception as e:
                    print(f"Error converting table to dataframe: {e}")

            elif "chart" in shape_type:
                df = chart_shape_to_df(shape)
                if df is not None and not df.empty:
                    try:
                        content_elements.append(df.to_dict(orient='records'))
                    except Exception as e:
                        print(f"Error converting table to dataframe: {e}")
        return content_elements

    def _extract_pptx_elements(self, slide) -> Dict[str, Any]:
        content_elements = []
        for idx, shape in enumerate(slide.shapes):
            shape_type = self._get_shape_type(shape)
            if shape_type == "text":
                content_elements.append({
                    "id": idx,
                    "type": 'textBox',
                    "text": shape.text.strip(),
                    'role': '',
                    "layout": get_shape_layout(shape)
                })
            if shape_type == "table":
                df = table_shape_to_df(shape)
                if df is not None and not df.empty:
                    content_elements.append({
                        "id": idx,
                        "type": 'table',
                        'role': '',
                        "layout": get_shape_layout(shape),
                        "data": df.to_dict(orient='records'),
                    })
            elif "chart" in shape_type:
                df = chart_shape_to_df(shape)
                if df is not None and not df.empty:
                    content_elements.append({
                        "id": idx,
                        "type": 'chart',
                        'role': '',
                        "layout": get_shape_layout(shape),
                        "data": df.to_dict(orient='records'),
                    })
        template_slide = {
            "slide_size": {"width": emu_to_cm(self.presentation.slide_width),
                           "height": emu_to_cm(self.presentation.slide_height)},
            "elements": content_elements
        }
        return template_slide

    def _get_vlm_analysis(self, slide_idx: int) -> Tuple[List[Dict[str, Any]], int, int]:
        output_dir = Path(".").resolve()
        img, img_w, img_h, temp_path = _convert_ppt_to_image(
            self.file_path, slide_idx + 1, str(output_dir)
        )
        if not img:
            return [], 0, 0

        buffered = io.BytesIO()
        img.save(buffered, format="JPEG")
        base64_image = base64.b64encode(buffered.getvalue()).decode("utf-8")

        vlm_results = _call_vision_model_v2(base64_image)
        if temp_path and os.path.exists(temp_path):
            shutil.rmtree(temp_path)

        return vlm_results, img_w, img_h

    def _bbox_px_2_layout(self, bbox: list, pptx_cm_to_px_w_ratio: float, pptx_cm_to_px_h_ratio: float) -> dict:
        return {
            "x": round(bbox[0] / pptx_cm_to_px_w_ratio, 2),
            "y": round(bbox[1] / pptx_cm_to_px_h_ratio, 2),
            "width": round((bbox[2] - bbox[0]) / pptx_cm_to_px_w_ratio, 2),
            "height": round((bbox[3] - bbox[1]) / pptx_cm_to_px_h_ratio, 2),
        }

    def _nearest_point(self, point, points):
        px, py = point
        best_dist = float('inf')
        best_idx = -1

        for i, (x, y) in enumerate(points):
            d = hypot(x - px, y - py)
            if d < best_dist:
                best_dist = d
                best_idx = i
        return best_idx

    def pptx_box_layout_to_box_layout(self, layout, pptx_cm_to_px_w_ratio, pptx_cm_to_px_h_ratio):
        x1 = layout.get('x') * pptx_cm_to_px_w_ratio
        y1 = layout.get('y') * pptx_cm_to_px_h_ratio
        x2 = x1 + layout.get('width') * pptx_cm_to_px_w_ratio
        y2 = y1 + layout.get('height') * pptx_cm_to_px_h_ratio
        return (x1, y1, x2, y2)

    def calculate_iou(self, box1: Tuple[float, float, float, float],
                      box2: Tuple[float, float, float, float]) -> float:
        x1, y1, w1, h1 = box1
        x2, y2, w2, h2 = box2

        box1_x2, box1_y2 = x1 + w1, y1 + h1
        box2_x2, box2_y2 = x2 + w2, y2 + h2

        inter_x1 = max(x1, x2)
        inter_y1 = max(y1, y2)
        inter_x2 = min(box1_x2, box2_x2)
        inter_y2 = min(box1_y2, box2_y2)

        inter_width = max(0, inter_x2 - inter_x1)
        inter_height = max(0, inter_y2 - inter_y1)
        inter_area = inter_width * inter_height

        box1_area = w1 * h1
        box2_area = w2 * h2
        union_area = box1_area + box2_area - inter_area

        if union_area == 0:
            return 0.0

        iou = inter_area / union_area
        return iou

    def _match_elements_with_iou(self, pptx_args: Dict[str, Any], vlm_results: List[Dict[str, Any]], img_w: int,
                                 img_h: int, threshold: float = 0.1) -> Dict[str, Any]:
        pptx_cm_to_px_w_ratio = img_w / emu_to_cm(self.presentation.slide_width)
        pptx_cm_to_px_h_ratio = img_h / emu_to_cm(self.presentation.slide_height)

        pptx_elements = pptx_args.get('elements', [])
        used_list2 = set()
        for i, item1 in enumerate(pptx_elements):
            box1 = self.pptx_box_layout_to_box_layout(item1['layout'], pptx_cm_to_px_w_ratio, pptx_cm_to_px_h_ratio)
            best_match_idx = -1
            best_iou = threshold
            for j, item2 in enumerate(vlm_results):
                if j in used_list2:
                    continue
                box2 = tuple(item2['bbox'])
                iou = self.calculate_iou(box1, box2)
                if iou > best_iou:
                    best_iou = iou
                    best_match_idx = j

            if best_match_idx != -1:
                pptx_elements[i]['role'] = vlm_results[best_match_idx]['shape_type']

                used_list2.add(best_match_idx)
            else:
                print(box1)
                print(item1)

        template_slide = {
            "slide_size": pptx_args.get('slide_size'),
            "elements": pptx_elements
        }
        return template_slide

    def _match_elements(self, pptx_elements: Dict[str, Any], vlm_results: List[Dict[str, Any]], img_w: int,
                        img_h: int) -> Dict[str, Any]:
        pptx_cm_to_px_w_ratio = img_w / emu_to_cm(self.presentation.slide_width)
        pptx_cm_to_px_h_ratio = img_h / emu_to_cm(self.presentation.slide_height)
        elements = pptx_elements.get('elements', [])
        points = []
        for element in elements:
            points.append((element.get('layout').get('x'), element.get('layout').get('y')))

        for item in vlm_results:
            item_point = (
            round(item["bbox"][0] / pptx_cm_to_px_w_ratio, 2), round(item["bbox"][1] / pptx_cm_to_px_h_ratio, 2))
            nearest_point_idx = self._nearest_point(item_point, points)
            elements[nearest_point_idx]['role'] = item['shape_type']

        template_slide = {
            "slide_size": pptx_elements.get('slide_size'),
            "elements": elements
        }

        return template_slide

    def _match_caption_and_table(self, structured_data) -> Dict[str, Any]:
        updated_elements = []
        elements = structured_data.get('elements', [])
        elements_table = [item for item in elements if
                          item.get('role') == 'table' or item.get('role') == 'chart-bar' or item.get(
                              'role') == 'chart-line']
        for item in elements:
            if item.get('role') == 'slide-title' or item.get('role') == 'body-text':
                updated_elements.append(item)
            if item.get('role') == 'caption':
                caption = item.get('text')
                caption_point = (item.get('layout').get('x'), item.get('layout').get('y'))
                table_points = []
                for elements_table_layout in elements_table:
                    table_points.append(
                        (elements_table_layout.get('layout').get('x'), elements_table_layout.get('layout').get('y')))
                best_idx = self._nearest_point(caption_point, table_points)
                table_data = elements_table[best_idx].get('data')

                chain = self.table_information_extraction_prompt_template | self.model
                try:
                    table_args = chain.invoke({"table_caption": caption, "table_data": table_data}).content
                except Exception as e:
                    print(e)
                table_args = ast.literal_eval(table_args)
                print(table_args)
                elements_table[best_idx]["args"] = table_args
                updated_elements.append(item)
                updated_elements.append(elements_table[best_idx])

        return {"slide_size": structured_data.get("slide_size"), "elements": updated_elements}

    def standardize_table_keys(self, table_data, target_key='category'):
        if not table_data or not isinstance(table_data, list):
            return table_data

        first_row = table_data[0]
        if target_key not in first_row:
            return table_data

        sample_val = str(first_row[target_key]).strip()

        new_key_name = target_key

        if re.match(r'^(19|20)\d{2}$', sample_val):
            new_key_name = 'year'

        elif re.match(r'^(19|20)\d{2}[-/.](0[1-9]|1[0-2])$', sample_val):
            new_key_name = 'month'

        elif re.search(r'(m²|㎡|sqm)', sample_val, re.IGNORECASE):
            new_key_name = 'area_range'

        elif re.search(r'(M|W|k|Dollar)', sample_val, re.IGNORECASE):
            new_key_name = 'price_range'

        if new_key_name != target_key:
            new_table_data = []
            for row in table_data:
                target_val = row.get(target_key)
                new_row = {new_key_name: target_val}
                for k, v in row.items():
                    if k != target_key:
                        new_row[k] = v

                new_table_data.append(new_row)
            return new_table_data

        return table_data

    def _match_caption_and_table1(self, structured_data, slide_filters):
        updated_elements = []
        updated_slide_filters = []
        num = 0
        elements = structured_data.get('elements', [])
        elements_table = [item for item in elements if
                          item.get('role') == 'table' or item.get('role') == 'chart-bar' or item.get(
                              'role') == 'chart-line']
        for item in elements:
            if item.get('role') == 'slide-title' or item.get('role') == 'body-text':
                updated_elements.append(item)
            if item.get('role') == 'caption':
                caption = item.get('text')
                caption_point = (item.get('layout').get('x'), item.get('layout').get('y'))
                table_points = []
                for elements_table_layout in elements_table:
                    table_points.append(
                        (elements_table_layout.get('layout').get('x'), elements_table_layout.get('layout').get('y')))
                best_idx = self._nearest_point(caption_point, table_points)
                table_data = elements_table[best_idx].get('data')
                table_data = self.standardize_table_keys(table_data)

                chain = self.table_information_extraction_prompt_template | self.model
                try:
                    table_args = chain.invoke({"table_caption": caption, "table_data": table_data}).content

                    table_args = re.sub(r'<think>.*?</think>', '', table_args, flags=re.DOTALL).strip()

                except Exception as e:
                    print(e)
                table_args = ast.literal_eval(table_args)
                slide_filters[num].setdefault('fun_tool', {})['quadruples'] = table_args
                updated_slide_filters.append(slide_filters[num])
                num+=1
                elements_table[best_idx]["args"] = table_args
                updated_elements.append(item)
                updated_elements.append(elements_table[best_idx])

        return updated_slide_filters
    def parse_slide_vlm(self, slide_idx: int = 0) -> Dict[str, Any]:

        if slide_idx >= len(self.presentation.slides):
            raise IndexError(f"Slide index {slide_idx} out of range.")
        slide = self.presentation.slides[slide_idx]
        pptx_elements = self._extract_pptx_elements(slide)
        vlm_results, img_w, img_h = self._get_vlm_analysis(slide_idx)
        if not vlm_results:
            return {"error": "Unable to generate slide image or call VLM"}

        structured_data = self._match_elements(pptx_elements, vlm_results, img_w, img_h)
        # structured_data = self._match_elements_with_iou(pptx_elements, vlm_results, img_w, img_h)

        data = structured_data
        template_slide = self._match_caption_and_table(data)
        return template_slide

    @staticmethod
    def save_dict_as_yaml(data: Dict, output_path: Path):
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            yaml.dump(data, f, default_flow_style=False, allow_unicode=True, indent=2, sort_keys=False)
        print(f"Successfully saved extracted structure to: {output_path}")

